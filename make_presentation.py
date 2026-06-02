# -*- coding: utf-8 -*-
import shutil, sys, os
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

EXAMPLE = r"C:\Users\thund\Downloads\Telegram Desktop\Пример_презентации.pptx"
OUTPUT  = r"C:\Users\thund\OneDrive\Рабочий стол\MyFitness_Презентация.pptx"
IMGS    = r"C:\Users\thund\OneDrive\Рабочий стол\docx_images"
DESK    = r"C:\Users\thund\OneDrive\Рабочий стол"

WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BLUE   = RGBColor(0x3B, 0x5B, 0xFF)
GRAY   = RGBColor(0xB0, 0xB8, 0xC1)
RED    = RGBColor(0xFF, 0x6B, 0x7A)
LGRAY  = RGBColor(0xD0, 0xD8, 0xE8)

EMU = 914400.0

shutil.copy(EXAMPLE, OUTPUT)
prs = Presentation(OUTPUT)

W = prs.slide_width.inches   # 13.33
H = prs.slide_height.inches  # 7.50

# ─── helpers ─────────────────────────────────────────────────────────────────

def set_run_text(para, text):
    """Set text in first run of paragraph, preserve run XML rPr."""
    runs = para.runs
    if runs:
        runs[0].text = text
        # Clear extra runs
        for r in runs[1:]:
            r.text = ''
    else:
        run = para.add_run()
        run.text = text

def set_ph(slide, idx, text):
    """Update placeholder text (first paragraph, first run)."""
    try:
        ph = slide.placeholders[idx]
    except KeyError:
        return
    tf = ph.text_frame
    paras = tf.paragraphs
    if paras:
        set_run_text(paras[0], text)
        # Remove extra paragraphs' text
        for p in paras[1:]:
            for r in p.runs:
                r.text = ''

def rm_textboxes(slide):
    """Remove all TEXT_BOX shapes from slide."""
    to_del = [s for s in slide.shapes if s.shape_type == 17]
    for s in to_del:
        s._element.getparent().remove(s._element)

def rm_pictures(slide):
    """Remove all PICTURE shapes from slide."""
    to_del = [s for s in slide.shapes if s.shape_type == 13]
    for s in to_del:
        s._element.getparent().remove(s._element)

def rm_autoshapes(slide):
    """Remove all non-placeholder AUTO_SHAPE (type=1) shapes from slide."""
    to_del = [s for s in slide.shapes
              if s.shape_type == 1 and not s.is_placeholder]
    for s in to_del:
        s._element.getparent().remove(s._element)

def clear_ph_text(slide, idx):
    """Set placeholder text to empty string."""
    try:
        ph = slide.placeholders[idx]
        for para in ph.text_frame.paragraphs:
            for run in para.runs:
                run.text = ''
    except KeyError:
        pass

def add_tb(slide, l, t, w, h, text,
           size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
           italic=False, wrap=True, name='Arial'):
    """Add a single-run text box."""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.color.rgb = color
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = name
    return tb

def add_paras(slide, l, t, w, h, items, size=16, color=WHITE, sp_after=5):
    """Add multi-paragraph text box.
    items: list of str or (str, bold) or (str, bold, RGBColor)
    """
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if isinstance(item, str):
            text, bold, col = item, False, color
        elif len(item) == 2:
            text, bold, col = item[0], item[1], color
        else:
            text, bold, col = item[0], item[1], item[2]
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(sp_after)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = col
        run.font.name = 'Arial'
    return tb

def add_img(slide, path, l, t, w, h):
    return slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))

# ─── Slide 1: Title ──────────────────────────────────────────────────────────
s1 = prs.slides[0]
rm_textboxes(s1)

set_ph(s1, 0, 'Курсовая работа')
set_ph(s1, 1,
    'Разработка мобильного Android-приложения\n'
    '«MyFitness» для учёта питания пользователя')
set_ph(s1, 10, '31.05.2026')

add_tb(s1, 6.04, 5.74, 6.38, 0.55,
       'Студент: Михеев Альберт',
       size=16, color=LGRAY)

# ─── Slide 2: Введение и Актуальность ───────────────────────────────────────
s2 = prs.slides[1]
rm_textboxes(s2)
rm_pictures(s2)

set_ph(s2, 0, 'Введение и Актуальность')
set_ph(s2, 10, '31.05.2026')

add_paras(s2, 0.50, 1.30, 6.20, 2.20, [
    ('Актуальность темы', True, BLUE),
    ('По данным ВОЗ, более 1,9 млрд взрослых имеют '
     'избыточный вес — питание является главным '
     'управляемым фактором здоровья.', False),
    ('', False),
    ('Рынок фитнес-приложений превысил $14 млрд '
     'в 2023 году и продолжает расти на 17% ежегодно '
     '(Allied Market Research).', False),
], size=17, sp_after=4)

add_paras(s2, 0.50, 3.60, 6.20, 2.80, [
    ('Цель проекта', True, BLUE),
    ('Разработать клиент-серверное Android-приложение '
     'для удобного ежедневного учёта калорий и БЖУ с '
     'синхронизацией данных через REST API.', False),
], size=17, sp_after=4)

add_paras(s2, 7.00, 1.30, 5.80, 5.00, [
    ('Потребности пользователей:', True, BLUE),
    ('•  Простой подсчёт калорий и БЖУ (Б/Ж/У)', False),
    ('•  Персонализация: вес, рост, цель питания', False),
    ('•  Работа без постоянного интернета (офлайн)', False),
    ('•  Синхронизация данных с сервером', False),
    ('•  Тёмный интерфейс, удобный для ежедневного использования', False),
    ('', False),
    ('Задачи проекта:', True, BLUE),
    ('•  Аутентификация пользователей (JWT)', False),
    ('•  Дневник питания с приёмами пищи', False),
    ('•  Автоматический расчёт нормы калорий', False),
    ('•  REST API сервер на Spring Boot + PostgreSQL', False),
    ('•  Локальное кеширование (Room/SQLite)', False),
], size=16, sp_after=3)

# ─── Slide 3: Обзор аналогов ─────────────────────────────────────────────────
s3 = prs.slides[2]
rm_textboxes(s3)
rm_pictures(s3)

set_ph(s3, 0, 'Обзор существующих аналогов')
set_ph(s3, 10, '31.05.2026')

def add_analog(slide, l, name, pros, cons, color_accent=BLUE):
    add_tb(slide, l, 1.25, 3.00, 0.50, name, size=19, bold=True, color=color_accent)
    add_paras(slide, l, 1.80, 3.00, 2.10, [
        ('Плюсы:', True, RGBColor(0x6E, 0xE7, 0xB7)),
    ] + [(f'+ {p}', False) for p in pros], size=15, sp_after=2)
    add_paras(slide, l, 3.95, 3.00, 1.60, [
        ('Минусы:', True, RED),
    ] + [(f'− {c}', False) for c in cons], size=15, sp_after=2)

add_analog(s3, 0.40, 'MyFitnessPal',
    ['Огромная БД продуктов', 'Интеграция с устройствами', 'Развитое сообщество'],
    ['Полный функционал — платно', 'Перегруженный интерфейс', 'Приложение на английском'])

add_analog(s3, 3.60, 'FatSecret',
    ['Полностью бесплатный', 'Есть поддержка русского языка', 'Сканер штрихкодов'],
    ['Устаревший интерфейс', 'Ограниченная аналитика', 'Нет офлайн-синхронизации'])

add_analog(s3, 6.80, 'Cronometer',
    ['Детальный учёт микронутриентов', 'Точность данных', 'Экспорт данных'],
    ['Сложный для новичков', 'Интерфейс перегружен', 'Ограниченная локализация'])

# Наше приложение
add_tb(s3, 10.00, 1.25, 2.90, 0.50, '«MyFitness»', size=19, bold=True, color=BLUE)
add_paras(s3, 10.00, 1.80, 2.90, 3.70, [
    ('Преимущества:', True, RGBColor(0x6E, 0xE7, 0xB7)),
    ('✓  Тёмный минималистичный интерфейс', False),
    ('✓  Офлайн-работа (Room cache)', False),
    ('✓  Синхронизация с сервером', False),
    ('✓  Расчёт нормы калорий по профилю', False),
    ('✓  Три приёма пищи + разбивка БЖУ', False),
    ('✓  JWT-аутентификация', False),
    ('✓  Clean Architecture', False),
], size=15, sp_after=2)

# ─── Slide 4: Функциональные требования ──────────────────────────────────────
s4 = prs.slides[3]
rm_textboxes(s4)
rm_pictures(s4)

set_ph(s4, 0, 'Функциональные требования')
set_ph(s4, 10, '31.05.2026')

add_img(s4, os.path.join(IMGS, 'img_05.png'), 0.40, 1.20, 6.00, 5.40)

add_paras(s4, 6.70, 1.20, 6.20, 5.50, [
    ('Актор: Гость', True, BLUE),
    ('•  Войти в аккаунт (email + пароль)', False),
    ('•  Зарегистрироваться', False),
    ('', False),
    ('Актор: Пользователь', True, BLUE),
    ('•  Просмотр дневника питания', False),
    ('•  Навигация по дням (календарь)', False),
    ('•  Добавить / удалить / редактировать продукт', False),
    ('•  Выбрать тип приёма пищи', False),
    ('•  Просмотр статистики калорий и БЖУ', False),
    ('•  Редактировать профиль (вес, рост, цель)', False),
    ('•  Выбрать цель питания', False),
    ('•  Выйти из аккаунта', False),
], size=16, sp_after=3)

# ─── Slide 5: Технологии разработки ─────────────────────────────────────────
s5 = prs.slides[4]
rm_textboxes(s5)
rm_pictures(s5)

# Clear body placeholder if exists
try:
    set_ph(s5, 1, '')
except:
    pass

set_ph(s5, 0, 'Технологии разработки')
set_ph(s5, 10, '31.05.2026')

add_tb(s5, 0.50, 1.20, 5.90, 0.55, 'Клиентская часть (Android)', size=20, bold=True, color=BLUE)
add_paras(s5, 0.50, 1.85, 5.90, 4.60, [
    ('Kotlin', True, WHITE),
    ('  Основной язык разработки Android-приложения', False, GRAY),
    ('Jetpack Compose', True, WHITE),
    ('  Декларативный UI-фреймворк от Google', False, GRAY),
    ('Room (SQLite)', True, WHITE),
    ('  Локальное кеширование данных питания', False, GRAY),
    ('Retrofit2 + OkHttp', True, WHITE),
    ('  HTTP-клиент для работы с REST API', False, GRAY),
    ('Dagger 2', True, WHITE),
    ('  Внедрение зависимостей (DI)', False, GRAY),
    ('Kotlin Coroutines / Flow', True, WHITE),
    ('  Асинхронная работа и реактивные потоки данных', False, GRAY),
], size=16, sp_after=1)

add_tb(s5, 6.90, 1.20, 6.00, 0.55, 'Серверная часть (Backend)', size=20, bold=True, color=BLUE)
add_paras(s5, 6.90, 1.85, 6.00, 4.60, [
    ('Kotlin', True, WHITE),
    ('  Язык серверной разработки', False, GRAY),
    ('Spring Boot', True, WHITE),
    ('  Фреймворк для создания REST API', False, GRAY),
    ('Spring Security + JWT', True, WHITE),
    ('  Аутентификация и авторизация запросов', False, GRAY),
    ('Spring Data JPA (Hibernate)', True, WHITE),
    ('  ORM для работы с базой данных', False, GRAY),
    ('PostgreSQL', True, WHITE),
    ('  Реляционная СУБД для хранения данных', False, GRAY),
    ('IntelliJ IDEA', True, WHITE),
    ('  Среда разработки серверной части', False, GRAY),
], size=16, sp_after=1)

# Vertical divider line (thin rectangle)
div = s5.shapes.add_shape(1,  # MSO_SHAPE_TYPE.RECTANGLE=1
    Inches(6.60), Inches(1.20), Inches(0.03), Inches(5.30))
div.fill.solid()
div.fill.fore_color.rgb = RGBColor(0x3B, 0x5B, 0xFF)
div.line.fill.background()

# ─── Slide 6: Архитектура ────────────────────────────────────────────────────
s6 = prs.slides[5]
rm_textboxes(s6)
rm_pictures(s6)

set_ph(s6, 0, 'Архитектура клиентской части')
set_ph(s6, 10, '31.05.2026')

add_paras(s6, 0.50, 1.25, 5.50, 5.30, [
    ('Clean Architecture', True, BLUE),
    ('Код разделён на 3 независимых слоя,', False),
    ('что обеспечивает тестируемость и масштабируемость.', False),
    ('', False),
    ('Presentation Layer', True, WHITE),
    ('  ViewModel (MVVM) + Jetpack Compose UI', False, GRAY),
    ('  Зависит только от UseCase, не знает про сеть/БД', False, GRAY),
    ('', False),
    ('Domain Layer', True, WHITE),
    ('  UseCase, Repository (интерфейсы), Domain Models', False, GRAY),
    ('  Чистая бизнес-логика без Android-зависимостей', False, GRAY),
    ('', False),
    ('Data Layer', True, WHITE),
    ('  RepositoryImpl, DataSource (Room + Retrofit)', False, GRAY),
    ('  Offline-first: сначала кеш, затем сервер', False, GRAY),
    ('', False),
    ('DI: Dagger 2 — связывает все слои через AppModule', False, GRAY),
], size=16, sp_after=1)

# Architecture diagram as styled text
arch_box = s6.shapes.add_textbox(
    Inches(6.20), Inches(1.25), Inches(6.70), Inches(5.30))
tf = arch_box.text_frame
tf.word_wrap = False

layers = [
    ('┌──────────────────────────────────────────────┐', WHITE, 14, False),
    ('│         PRESENTATION  (UI + ViewModel)       │', BLUE,  15, True),
    ('│   HomeViewModel  AuthViewModel  ProfileVM    │', LGRAY, 13, False),
    ('│   Jetpack Compose (HomeScreen, AuthScreen…)  │', LGRAY, 13, False),
    ('├──────────────────────────────────────────────┤', WHITE, 14, False),
    ('│              DOMAIN  (UseCase)               │', RGBColor(0xA0, 0xC0, 0xFF), 15, True),
    ('│  LoadDayUseCase  AddFoodItemUseCase          │', LGRAY, 13, False),
    ('│  LoginUseCase  GetCaloriesGoalUseCase  …     │', LGRAY, 13, False),
    ('│  Repository interfaces (FoodRepo, UserRepo)  │', LGRAY, 13, False),
    ('├──────────────────────────────────────────────┤', WHITE, 14, False),
    ('│              DATA  (Repository + DS)         │', RGBColor(0xFF, 0xC0, 0x80), 15, True),
    ('│  FoodRepositoryImpl   UserRepositoryImpl     │', LGRAY, 13, False),
    ('│  Room (RoomStorageImpl)     Retrofit API     │', LGRAY, 13, False),
    ('│  RemoteFoodDataSource   RemoteUserDataSource │', LGRAY, 13, False),
    ('└──────────────────────────────────────────────┘', WHITE, 14, False),
]
first = True
for text, col, sz, bd in layers:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    p.space_after = Pt(1)
    run = p.add_run()
    run.text = text
    run.font.color.rgb = col
    run.font.size = Pt(sz)
    run.font.bold = bd
    run.font.name = 'Courier New'

# ─── Slide 7: Модель хранения данных ─────────────────────────────────────────
s7 = prs.slides[6]
rm_textboxes(s7)
rm_pictures(s7)

set_ph(s7, 0, 'Модель хранения данных')
set_ph(s7, 10, '31.05.2026')

add_img(s7, os.path.join(IMGS, 'img_15.png'), 0.40, 1.25, 8.00, 3.20)

add_paras(s7, 0.40, 4.55, 7.80, 2.80, [
    ('Серверная база данных (PostgreSQL):', True, BLUE),
    ('USERS: firebase_uid (PK), name, gender, email, weight, height, target, password_hash', False, LGRAY),
    ('DAY_FOOD: id (PK), firebase_uid (FK), date, calories, protein, fats, carbohydrates', False, LGRAY),
    ('FOOD_ITEMS: id (PK), day_food_id (FK), name, weight, type_of_meal, calories, protein, fats, carbohydrates', False, LGRAY),
], size=15, sp_after=3)

add_paras(s7, 8.60, 1.25, 4.50, 5.80, [
    ('Локальное хранилище (Room)', True, BLUE),
    ('', False),
    ('Таблица day_food_items:', True, WHITE),
    ('•  Зеркало серверных данных', False, GRAY),
    ('•  Offline-first чтение', False, GRAY),
    ('•  Обновляется при синхронизации', False, GRAY),
    ('', False),
    ('Таблица users:', True, WHITE),
    ('•  Кеш профиля пользователя', False, GRAY),
    ('•  Позволяет работать без сети', False, GRAY),
    ('', False),
    ('Стратегия синхронизации:', True, BLUE),
    ('1. Чтение → Room (мгновенно)', False, GRAY),
    ('2. Фоновый запрос к серверу', False, GRAY),
    ('3. Обновление Room и UI', False, GRAY),
    ('4. Запись → сначала локально,', False, GRAY),
    ('   затем синхронизация с API', False, GRAY),
], size=15, sp_after=2)

# ─── Slide 8: Экраны приложения ──────────────────────────────────────────────
s8 = prs.slides[7]
rm_textboxes(s8)
rm_pictures(s8)

set_ph(s8, 0, 'Экраны приложения')
set_ph(s8, 10, '31.05.2026')

# Full auth+main screens image (wide)
add_img(s8, os.path.join(IMGS, 'img_00.png'), 0.30, 1.20, 8.40, 5.50)

# Profile screen
add_img(s8, os.path.join(IMGS, 'img_02.png'), 9.00, 1.20, 4.00, 5.20)

# Caption labels
add_tb(s8, 0.40, 6.75, 4.00, 0.45, 'Экран входа', size=14, color=GRAY, align=PP_ALIGN.CENTER)
add_tb(s8, 3.70, 6.75, 5.00, 0.45, 'Главный экран — дневник питания', size=14, color=GRAY, align=PP_ALIGN.CENTER)
add_tb(s8, 9.00, 6.45, 4.00, 0.45, 'Экран профиля', size=14, color=GRAY, align=PP_ALIGN.CENTER)

# ─── Slide 9: Заключение ─────────────────────────────────────────────────────
s9 = prs.slides[8]
rm_textboxes(s9)
rm_pictures(s9)
rm_autoshapes(s9)
# Clear the "Демо-видео" placeholder (idx varies — clear all non-title/date/num placeholders)
for ph in s9.placeholders:
    idx = ph.placeholder_format.idx
    if idx not in (0, 10, 12):
        clear_ph_text(s9, idx)

set_ph(s9, 0, 'Заключение')
set_ph(s9, 10, '31.05.2026')

add_paras(s9, 0.50, 1.25, 6.20, 4.50, [
    ('Результаты работы', True, BLUE),
    ('', False),
    ('✓  Разработан Android-клиент на Kotlin + Jetpack Compose', False),
    ('✓  Реализована чистая архитектура (Clean Architecture + MVVM)', False),
    ('✓  Создан REST API-сервер на Kotlin + Spring Boot', False),
    ('✓  JWT-аутентификация (регистрация / вход)', False),
    ('✓  Пищевой дневник: добавление, редактирование, удаление', False),
    ('✓  Автоматический расчёт нормы калорий по профилю', False),
    ('✓  Offline-first: кеш Room + синхронизация с PostgreSQL', False),
    ('✓  Прогресс-бар калорий и цветные индикаторы БЖУ', False),
    ('✓  Профиль пользователя (пол, вес, рост, цель)', False),
], size=16, sp_after=3)

add_paras(s9, 7.10, 1.25, 5.70, 4.50, [
    ('Планы развития', True, BLUE),
    ('', False),
    ('→  База данных продуктов питания со сканером штрихкода', False),
    ('→  Экран статистики с графиками по неделям', False),
    ('→  Уведомления-напоминания о приёмах пищи', False),
    ('→  Интеграция с фитнес-трекерами', False),
    ('→  Водный баланс (трекер воды)', False),
    ('→  Социальные функции: профиль и прогресс', False),
], size=16, sp_after=4)

add_tb(s9, 0.50, 5.90, 12.30, 0.60,
    'В ходе выполнения курсовой работы разработано клиент-серверное приложение для учёта '
    'питания с использованием современных технологий Android-разработки и Java/Kotlin backend.',
    size=15, color=GRAY, wrap=True)

prs.save(OUTPUT)
print(f"\n✓ Презентация сохранена: {OUTPUT}")
